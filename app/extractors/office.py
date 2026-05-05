"""Office document extractor (.docx, .xlsx, .pptx).

Extracts text content from Office Open XML formats using pure Python
libraries (python-docx, openpyxl, python-pptx). No LibreOffice required.
"""

import logging
from pathlib import Path

from app.extractors.base import ContentExtractor, ExtractionResult, TextChunk

logger = logging.getLogger(__name__)

MAX_DOCX_BYTES = 50 * 1024 * 1024
MAX_XLSX_BYTES = 50 * 1024 * 1024
MAX_PPTX_BYTES = 100 * 1024 * 1024
MAX_XLSX_ROWS_PER_SHEET = 1000

EXTRACTOR_NAME = "office"


class OfficeExtractor(ContentExtractor):
    """Extracts text from .docx, .xlsx, and .pptx files."""

    supported_extensions: list[str] = [".docx", ".xlsx", ".pptx"]

    def extract(self, file_path: str) -> ExtractionResult:
        ext = Path(file_path).suffix.lower()
        empty = ExtractionResult(chunks=[], markdown=None, extractor=EXTRACTOR_NAME)

        try:
            if ext == ".docx":
                return self._extract_docx(file_path, empty)
            if ext == ".xlsx":
                return self._extract_xlsx(file_path, empty)
            if ext == ".pptx":
                return self._extract_pptx(file_path, empty)
        except Exception as e:
            logger.error("Office extraction failed for %s: %s", file_path, e)

        return ExtractionResult(chunks=[], markdown=None, extractor="office_error")

    def _extract_docx(self, file_path: str, empty: ExtractionResult) -> ExtractionResult:
        size = Path(file_path).stat().st_size
        if size > MAX_DOCX_BYTES:
            logger.warning("Skipping .docx (size %d > %d): %s", size, MAX_DOCX_BYTES, file_path)
            return empty

        import docx

        doc = docx.Document(file_path)
        sections: list[str] = []
        current_heading = ""
        current_lines: list[str] = []

        def flush() -> None:
            block = "\n".join(current_lines).strip()
            if block:
                sections.append(block)
            current_lines.clear()

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style_name = para.style.name if para.style else ""
            if style_name.startswith("Heading"):
                flush()
                try:
                    level = int(style_name.split()[-1])
                except ValueError:
                    level = 1
                prefix = "#" * min(level, 6)
                current_heading = f"{prefix} {text}"
                current_lines.append(current_heading)
            else:
                current_lines.append(text)

        for table in doc.tables:
            flush()
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    current_lines.append(" | ".join(cells))

        flush()

        full_text = "\n\n".join(sections)
        if not full_text.strip():
            return empty

        raw_chunks = self.chunk_text(full_text)
        chunks = [
            TextChunk(
                text=c,
                metadata=f"section: {current_heading}" if current_heading else "",
            )
            for c in raw_chunks
        ]
        return ExtractionResult(chunks=chunks, markdown=None, extractor=EXTRACTOR_NAME)

    def _extract_xlsx(self, file_path: str, empty: ExtractionResult) -> ExtractionResult:
        size = Path(file_path).stat().st_size
        if size > MAX_XLSX_BYTES:
            logger.warning("Skipping .xlsx (size %d > %d): %s", size, MAX_XLSX_BYTES, file_path)
            return empty

        import openpyxl

        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        all_chunks: list[TextChunk] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text: list[str] = []
            row_num = 0

            for row in ws.iter_rows(values_only=True):
                if row_num >= MAX_XLSX_ROWS_PER_SHEET:
                    break
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    rows_text.append(" ".join(cells))
                    row_num += 1

            if not rows_text:
                continue

            section_text = f"## Sheet: {sheet_name}\n\n" + "\n".join(rows_text)
            raw_chunks = self.chunk_text(section_text)
            truncated = row_num >= MAX_XLSX_ROWS_PER_SHEET
            suffix = ", truncated" if truncated else ""

            for i, c in enumerate(raw_chunks):
                start = i * (len(rows_text) // max(len(raw_chunks), 1))
                end = min(start + len(rows_text) // max(len(raw_chunks), 1), row_num)
                all_chunks.append(
                    TextChunk(
                        text=c,
                        metadata=f"sheet: {sheet_name}, rows: {start}-{end}{suffix}",
                    )
                )

        wb.close()
        if not all_chunks:
            return empty
        return ExtractionResult(chunks=all_chunks, markdown=None, extractor=EXTRACTOR_NAME)

    def _extract_pptx(self, file_path: str, empty: ExtractionResult) -> ExtractionResult:
        size = Path(file_path).stat().st_size
        if size > MAX_PPTX_BYTES:
            logger.warning("Skipping .pptx (size %d > %d): %s", size, MAX_PPTX_BYTES, file_path)
            return empty

        from pptx import Presentation

        prs = Presentation(file_path)
        all_chunks: list[TextChunk] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            lines: list[str] = []

            title_shape = slide.shapes.title
            title_text = title_shape.text.strip() if title_shape and title_shape.has_text_frame else ""

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if shape is title_shape:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)

            if not title_text and not lines:
                continue

            heading = f"## Slide {slide_num}: {title_text}" if title_text else f"## Slide {slide_num}"
            slide_text = heading + ("\n\n" + "\n".join(lines) if lines else "")
            raw_chunks = self.chunk_text(slide_text)

            for c in raw_chunks:
                all_chunks.append(TextChunk(text=c, metadata=f"slide: {slide_num}"))

        if not all_chunks:
            return empty
        return ExtractionResult(chunks=all_chunks, markdown=None, extractor=EXTRACTOR_NAME)
